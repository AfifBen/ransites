from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'output'
ASSETS = ROOT / 'presentation_assets'
OUT.mkdir(parents=True, exist_ok=True)


def add_title(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(18, 33, 52)

    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(12), Inches(2.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(187, 210, 238)

    date_box = s.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(8), Inches(0.8))
    d = date_box.text_frame.paragraphs[0]
    d.text = f'Presentation technique - {datetime.now().strftime("%d/%m/%Y")}'
    d.font.size = Pt(15)
    d.font.color.rgb = RGBColor(220, 230, 245)


def add_bullets(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)

    b = s.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.8), Inches(5.7))
    tfb = b.text_frame
    tfb.clear()
    for i, line in enumerate(bullets):
        pp = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        pp.text = line
        pp.font.size = Pt(21)
        pp.font.color.rgb = RGBColor(45, 55, 75)


def add_image(prs, title, img_path, caption):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)

    if img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.2), width=Inches(11.7))
    else:
        ph = s.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.8))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(238, 243, 250)
        ph.line.color.rgb = RGBColor(200, 210, 225)

    c = s.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    cp = c.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(13)
    cp.alignment = PP_ALIGN.CENTER


def add_two_images(prs, title, left, lcap, right, rcap):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)

    if left.exists():
        s.shapes.add_picture(str(left), Inches(0.7), Inches(1.2), width=Inches(5.8))
    if right.exists():
        s.shapes.add_picture(str(right), Inches(6.8), Inches(1.2), width=Inches(5.8))

    c1 = s.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(5.8), Inches(0.4))
    c1.text_frame.paragraphs[0].text = lcap
    c1.text_frame.paragraphs[0].font.size = Pt(13)
    c1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    c2 = s.shapes.add_textbox(Inches(6.8), Inches(6.7), Inches(5.8), Inches(0.4))
    c2.text_frame.paragraphs[0].text = rcap
    c2.text_frame.paragraphs[0].font.size = Pt(13)
    c2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_lifecycle(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    p.text = 'Version suivante - Cycle de vie D1 -> On Air'
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)

    steps = [
        ('D1', 'Besoin RAN', RGBColor(31, 78, 121)),
        ('D2', 'Acquisition', RGBColor(46, 117, 182)),
        ('D3', 'Design technique', RGBColor(68, 114, 196)),
        ('D4', 'Construction', RGBColor(112, 173, 71)),
        ('D5', 'Integration / Tests', RGBColor(255, 192, 0)),
        ('On Air', 'Mise en service', RGBColor(192, 0, 0)),
    ]

    x = 0.7
    y = 2.0
    for i, (code, label, color) in enumerate(steps):
        box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.9), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)

        tb = s.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.2), Inches(1.65), Inches(1.6))
        tf = tb.text_frame
        pp = tf.paragraphs[0]
        pp.text = code
        pp.font.bold = True
        pp.font.size = Pt(20)
        pp.font.color.rgb = RGBColor(255, 255, 255)

        pp2 = tf.add_paragraph()
        pp2.text = label
        pp2.font.size = Pt(12)
        pp2.font.color.rgb = RGBColor(245, 245, 245)

        if i < len(steps) - 1:
            arr = s.shapes.add_textbox(Inches(x + 1.95), Inches(y + 0.8), Inches(0.35), Inches(0.5))
            ap = arr.text_frame.paragraphs[0]
            ap.text = '>'
            ap.font.size = Pt(22)
            ap.font.bold = True
            ap.font.color.rgb = RGBColor(90, 100, 120)

        x += 2.1

    note = s.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(12.0), Inches(2.1))
    nt = note.text_frame
    for i, line in enumerate([
        'Objectif VNext: orchestrer les interactions RAN / Acquisition / Construction dans un seul workflow.',
        'Chaque jalon sera trace (statut, responsable, date cible/reelle, blocages, evidences documentaires).',
        'Resultat attendu: delais reduits, meilleure coordination transverse et go-live plus predictible.'
    ]):
        p = nt.paragraphs[0] if i == 0 else nt.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(45, 55, 75)


def generate():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    dash = ASSETS / 'real_dashboard.png'
    sites = ASSETS / 'real_sites_table.png'
    profile = ASSETS / 'real_site_profile.png'
    imp = ASSETS / 'real_import_export.png'
    c1 = ASSETS / 'chart_entities.png'
    c2 = ASSETS / 'chart_quality.png'

    add_title(prs, 'RANSites', 'Plateforme simple et technique pour le departement RAN')

    add_bullets(prs, 'Besoin du departement RAN', [
        'Consolider les donnees Sites / Sectors / Cells dans un referentiel unique.',
        'Reduire les erreurs de mapping et accelerer les livrables planning.',
        'Standardiser les flux Import/Export (Allplan, KML) avec controle qualite.',
        'Maitriser les acces par perimetre geographique (wilaya/commune/site).',
    ])

    add_two_images(prs, 'Application en action', dash, 'Dashboard operationnel', sites, 'Table Sites & selection')
    add_two_images(prs, 'Valeur technique immediate', profile, 'Site Profile: KPI + voisins + beams', imp, 'Import/Export + templates')

    add_two_images(prs, 'Indicateurs de pilotage', c1, 'Couverture des donnees', c2, 'Qualite des donnees')

    add_bullets(prs, 'Pourquoi adopter RANSites maintenant', [
        'Gain de temps operationnel sur les consolidations et exports recurrents.',
        'Qualite de donnees plus stable grace aux validations et profils techno.',
        'Visibilite management via KPI de qualite et couverture.',
        'Base solide pour industrialiser le cycle de vie complet des sites.',
    ])

    add_lifecycle(prs)

    add_bullets(prs, 'Decision proposee', [
        'Valider RANSites comme outil de reference du departement RAN.',
        'Lancer la VNext cycle de vie D1 -> On Air avec Acquisition et Construction.',
        'Nommer un sponsor metier et un owner technique pour la generalisation.'
    ])

    out = OUT / 'RANSites_Presentation_Simple_Technique.pptx'
    prs.save(out)
    print(out)


if __name__ == '__main__':
    generate()
