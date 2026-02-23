import os
import sqlite3
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DB_PATH = ROOT / 'instance' / 'radio.db'
ASSETS = ROOT / 'presentation_assets'
OUT = ROOT / 'output'
ASSETS.mkdir(parents=True, exist_ok=True)
OUT.mkdir(parents=True, exist_ok=True)


def q1(conn, sql, params=()):
    cur = conn.cursor()
    cur.execute(sql, params)
    row = cur.fetchone()
    return row[0] if row else 0


def qall(conn, sql, params=()):
    cur = conn.cursor()
    cur.execute(sql, params)
    return cur.fetchall()


def safe_count(conn, table):
    return q1(conn, f"SELECT COUNT(*) FROM {table}")


def build_stats():
    conn = sqlite3.connect(DB_PATH)
    stats = {}

    stats['sites'] = safe_count(conn, 'site')
    stats['sectors'] = safe_count(conn, 'sector')
    stats['cells'] = safe_count(conn, 'cell')
    stats['regions'] = safe_count(conn, 'region')
    stats['wilayas'] = safe_count(conn, 'wilaya')
    stats['communes'] = safe_count(conn, 'commune')
    stats['suppliers'] = safe_count(conn, 'supplier')
    stats['antennas'] = safe_count(conn, 'antenna')
    stats['mapping'] = safe_count(conn, 'mapping')

    tech_rows = qall(conn, "SELECT UPPER(COALESCE(technology, 'N/A')), COUNT(*) FROM cell GROUP BY UPPER(COALESCE(technology, 'N/A'))")
    tech_counts = {'2G': 0, '3G': 0, '4G': 0, '5G': 0, 'OTHER': 0}
    for tech, count in tech_rows:
        t = str(tech or '').strip().upper()
        if t in tech_counts:
            tech_counts[t] += count
        elif t and t != 'N/A':
            tech_counts['OTHER'] += count

    stats['tech_counts'] = tech_counts

    stats['sites_without_sectors'] = q1(conn, """
        SELECT COUNT(*)
        FROM site s
        LEFT JOIN sector sec ON sec.site_id = s.id
        WHERE sec.id IS NULL
    """)

    stats['sectors_without_cells'] = q1(conn, """
        SELECT COUNT(*)
        FROM sector sec
        LEFT JOIN cell c ON c.sector_id = sec.id
        WHERE c.id IS NULL
    """)

    stats['cells_without_sector'] = q1(conn, "SELECT COUNT(*) FROM cell WHERE sector_id IS NULL")
    stats['cells_without_antenna'] = q1(conn, "SELECT COUNT(*) FROM cell WHERE antenna_id IS NULL")

    mapped_codes = {str(r[0]).strip() for r in qall(conn, "SELECT DISTINCT cell_code FROM mapping") if r and r[0] is not None}
    cellnames = [str(r[0]).strip() for r in qall(conn, "SELECT cellname FROM cell") if r and r[0] is not None]

    def extract_code(cellname):
        if '_' not in cellname:
            return None
        return cellname.rsplit('_', 1)[-1].strip() or None

    mapped_cells = sum(1 for c in cellnames if extract_code(c) in mapped_codes)
    stats['mapped_cells'] = mapped_cells
    stats['mapping_coverage'] = round((mapped_cells / stats['cells'] * 100.0), 1) if stats['cells'] else 0.0

    conn.close()
    return stats


def save_entity_bar(stats):
    labels = ['Sites', 'Sectors', 'Cells', 'Mapping', 'Antennas']
    values = [stats['sites'], stats['sectors'], stats['cells'], stats['mapping'], stats['antennas']]

    plt.figure(figsize=(8, 4.2))
    bars = plt.bar(labels, values, color=['#1f4e79', '#2e75b6', '#5b9bd5', '#70ad47', '#ffc000'])
    plt.title('RANSites - Couverture des donnees')
    plt.ylabel('Nombre')
    plt.grid(axis='y', linestyle='--', alpha=0.3)
    for b in bars:
        y = b.get_height()
        plt.text(b.get_x() + b.get_width() / 2, y, f'{int(y)}', ha='center', va='bottom', fontsize=9)
    plt.tight_layout()
    path = ASSETS / 'chart_entities.png'
    plt.savefig(path, dpi=170)
    plt.close()
    return path


def save_tech_donut(stats):
    tech = stats['tech_counts']
    labels = []
    values = []
    for k in ['2G', '3G', '4G', '5G', 'OTHER']:
        if tech[k] > 0:
            labels.append(k)
            values.append(tech[k])

    if not values:
        labels, values = ['N/A'], [1]

    colors = ['#70ad47', '#5b9bd5', '#ffc000', '#c00000', '#7f7f7f'][:len(values)]
    plt.figure(figsize=(5.8, 4.4))
    wedges, texts, autotexts = plt.pie(values, labels=labels, autopct='%1.0f%%', startangle=130, colors=colors, pctdistance=0.8)
    centre_circle = plt.Circle((0, 0), 0.54, fc='white')
    fig = plt.gcf()
    fig.gca().add_artist(centre_circle)
    plt.title('Distribution technologique des cellules')
    plt.tight_layout()
    path = ASSETS / 'chart_tech.png'
    plt.savefig(path, dpi=170)
    plt.close()
    return path


def save_quality_chart(stats):
    labels = ['Sites sans sectors', 'Sectors sans cells', 'Cells sans sector', 'Cells sans antenna']
    values = [
        stats['sites_without_sectors'],
        stats['sectors_without_cells'],
        stats['cells_without_sector'],
        stats['cells_without_antenna'],
    ]
    plt.figure(figsize=(8, 4.2))
    bars = plt.barh(labels, values, color=['#ed7d31', '#a5a5a5', '#4472c4', '#c55a11'])
    plt.title('Indicateurs qualite des donnees')
    plt.xlabel('Nombre')
    for b in bars:
        x = b.get_width()
        plt.text(x + 0.3, b.get_y() + b.get_height() / 2, f'{int(x)}', va='center', fontsize=9)
    plt.tight_layout()
    path = ASSETS / 'chart_quality.png'
    plt.savefig(path, dpi=170)
    plt.close()
    return path


def make_mock_screenshot(title, subtitle, filename, color=(31, 78, 121)):
    img = Image.new('RGB', (1600, 900), (239, 242, 247))
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, 1600, 86), fill=color)
    d.rectangle((0, 86, 280, 900), fill=(38, 48, 66))
    d.rounded_rectangle((320, 130, 1540, 840), radius=18, fill=(255, 255, 255), outline=(205, 210, 220), width=2)

    f1 = ImageFont.load_default()
    d.text((24, 30), f'RANSites - {title}', fill=(255, 255, 255), font=f1)
    d.text((345, 165), subtitle, fill=(80, 90, 110), font=f1)

    for i in range(5):
        y = 220 + i * 95
        d.rounded_rectangle((350, y, 1510, y + 66), radius=10, fill=(247, 249, 253), outline=(228, 232, 240), width=1)

    path = ASSETS / filename
    img.save(path)
    return path


def add_title_slide(prs, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(20, 35, 55)

    box = s.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11.5), Inches(2.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = 'RANSites'
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = 'Plateforme de gouvernance et planification RAN'
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(180, 205, 235)

    box2 = s.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(11), Inches(1.8))
    tf2 = box2.text_frame
    t = tf2.paragraphs[0]
    t.text = subtitle
    t.font.size = Pt(20)
    t.font.color.rgb = RGBColor(255, 255, 255)


def add_bullets_slide(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(31, 78, 121)

    b = s.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.6), Inches(5.6))
    tfb = b.text_frame
    tfb.clear()
    for i, line in enumerate(bullets):
        pp = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        pp.text = line
        pp.level = 0
        pp.font.size = Pt(22)
        pp.font.color.rgb = RGBColor(50, 60, 80)


def add_two_images_slide(prs, title, img1, cap1, img2, cap2):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.9))
    tf = t.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)

    s.shapes.add_picture(str(img1), Inches(0.7), Inches(1.3), width=Inches(5.9))
    s.shapes.add_picture(str(img2), Inches(6.7), Inches(1.3), width=Inches(5.9))

    c1 = s.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(5.9), Inches(0.4))
    c1.text_frame.paragraphs[0].text = cap1
    c1.text_frame.paragraphs[0].font.size = Pt(14)
    c1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    c2 = s.shapes.add_textbox(Inches(6.7), Inches(6.7), Inches(5.9), Inches(0.4))
    c2.text_frame.paragraphs[0].text = cap2
    c2.text_frame.paragraphs[0].font.size = Pt(14)
    c2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_chart_slide(prs, title, chart_path, side_bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.9))
    tf = t.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)

    s.shapes.add_picture(str(chart_path), Inches(0.7), Inches(1.2), width=Inches(7.3))

    b = s.shapes.add_textbox(Inches(8.2), Inches(1.4), Inches(4.1), Inches(5.6))
    tfb = b.text_frame
    tfb.clear()
    for i, line in enumerate(side_bullets):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(55, 65, 85)


def add_roadmap_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.9))
    tf = t.text_frame
    tf.paragraphs[0].text = 'Plan de deploiement (90 jours)'
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)

    phases = [
        ('S1-S2', 'Pilotage region', 'Onboarding, import baseline, formation'),
        ('S3-S6', 'Industrialisation', 'Generalisation progressive, support utilisateurs'),
        ('S7-S10', 'Extension nationale', 'Harmonisation process et KPIs adoption'),
        ('S11-S12', 'Bilan executif', 'ROI, governance finale, roadmap V2'),
    ]

    x = 0.8
    for i, (period, title, desc) in enumerate(phases):
        color = [(31, 78, 121), (46, 117, 182), (91, 155, 213), (112, 173, 71)][i]
        shape = s.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(2.8), Inches(4.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(255, 255, 255)

        tx = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.85), Inches(2.5), Inches(4.2))
        tf2 = tx.text_frame
        p0 = tf2.paragraphs[0]
        p0.text = period
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(255, 255, 255)

        p1 = tf2.add_paragraph()
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)

        p2 = tf2.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(240, 245, 255)
        x += 3.0


def generate():
    if not DB_PATH.exists():
        raise FileNotFoundError(f'Database not found: {DB_PATH}')

    stats = build_stats()

    c_entities = save_entity_bar(stats)
    c_tech = save_tech_donut(stats)
    c_quality = save_quality_chart(stats)

    real1 = ASSETS / 'real_dashboard.png'
    real2 = ASSETS / 'real_site_profile.png'
    real3 = ASSETS / 'real_import_export.png'

    ss1 = real1 if real1.exists() else make_mock_screenshot('Dashboard', 'Vue KPI, qualite de donnees, alertes', 'screen_dashboard.png')
    ss2 = real2 if real2.exists() else make_mock_screenshot('Site Profile', 'KPI, map, beams sector, voisins', 'screen_site_profile.png', color=(46, 117, 182))
    ss3 = real3 if real3.exists() else make_mock_screenshot('Import / Export', 'Templates, import multi-tech, exports Allplan/KML', 'screen_import_export.png', color=(91, 155, 213))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    today = datetime.now().strftime('%d/%m/%Y')
    add_title_slide(prs, f'Presentation direction - version {today}')

    add_bullets_slide(prs, 'Pourquoi RANSites', [
        'Unifier les donnees RAN: Sites, Sectors, Cells, Mapping',
        'Reduire les erreurs manuelles et les delais de production',
        'Standardiser les imports/exports metier (Allplan, KML)',
        'Gouverner les acces par perimetre geo (wilaya/commune/site)',
    ])

    add_two_images_slide(prs, 'Experience produit - vues principales', ss1, 'Dashboard executif', ss2, 'Site Profile avec beams')
    add_two_images_slide(prs, 'Flux operationnel', ss3, 'Import / Export & templates', ss2, 'Analyse locale par site')

    add_chart_slide(prs, 'Couverture des donnees (base actuelle)', c_entities, [
        f"Sites: {stats['sites']}",
        f"Sectors: {stats['sectors']}",
        f"Cells: {stats['cells']}",
        f"Mapping: {stats['mapping']}",
        f"Antennas: {stats['antennas']}",
    ])

    add_chart_slide(prs, 'Distribution technologique', c_tech, [
        f"2G: {stats['tech_counts']['2G']}",
        f"3G: {stats['tech_counts']['3G']}",
        f"4G: {stats['tech_counts']['4G']}",
        f"5G: {stats['tech_counts']['5G']}",
        f"Mapping coverage: {stats['mapping_coverage']}%",
    ])

    add_chart_slide(prs, 'Qualite des donnees - points de vigilance', c_quality, [
        f"Sites sans sectors: {stats['sites_without_sectors']}",
        f"Sectors sans cells: {stats['sectors_without_cells']}",
        f"Cells sans sector: {stats['cells_without_sector']}",
        f"Cells sans antenna: {stats['cells_without_antenna']}",
        'Actions: nettoyage cible + regles de validation',
    ])

    add_roadmap_slide(prs)

    add_bullets_slide(prs, 'Decision demandee', [
        'Valider un pilote officiel RANSites sur 1 region',
        'Nommer un sponsor metier + un owner technique',
        'Adopter RANSites comme referentiel de donnees RAN',
        'Planifier la generalisation entreprise apres pilote',
    ])

    out_file = OUT / 'RANSites_Presentation_Direction.pptx'
    prs.save(out_file)
    print(str(out_file))


if __name__ == '__main__':
    generate()
